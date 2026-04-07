from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

def create_pitch_deck():
    prs = Presentation()

    # --- STYLE ---
    # Helper for adding background color (simulated with a shape)
    def add_bg(slide, color_rgb):
        left = top = 0
        width = prs.slide_width
        height = prs.slide_height
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_rgb
        shape.line.width = 0

    # Colors
    BLUE_DARK = RGBColor(15, 23, 42)
    BLUE_ACCENT = RGBColor(59, 130, 246)
    GREEN_ACCENT = RGBColor(16, 185, 129)
    WHITE = RGBColor(255, 255, 255)

    # --- SLIDE 1: TITLE ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    add_bg(slide, BLUE_DARK)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "SmartAir Cameroon"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "IA pour la Résilience Climatique et Sanitaire\nPrédire l'invisible pour protéger le visible."
    subtitle.text_frame.paragraphs[0].font.color.rgb = BLUE_ACCENT

    # --- SLIDE 2: PROBLEM ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "L'enjeu de Santé Publique"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Le Défi Camerounais :"
    
    p = tf.add_paragraph()
    p.text = "• Impact de l'Harmattan et pics de chaleur saisonniers."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Hausse de 30% des maladies respiratoires en saison sèche."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Lacune : Manque de capteurs physiques PM2.5 coûteux."
    p.level = 1

    # --- SLIDE 3: SOLUTION ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "La Solution : SmartAir AI"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Transformer les données météo en bouclier sanitaire :"
    
    p = tf.add_paragraph()
    p.text = "• Modèle XGBoost haute performance (R² = 0.998)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Prédiction du proxy PM2.5 à partir de variables standards."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 16 features climatiques analysées (vent, temp, radiation)."
    p.level = 1

    # --- SLIDE 4: PRODUCT DEMO ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Un Outil de Décision Interactif"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Fonctionnalités Clés :"
    
    p = tf.add_paragraph()
    p.text = "• Carte de chaleur (Heatmap) dynamique du territoire."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Dashboard analytique Climat vs Pollution."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Simulateur de prévisions 'Live' pour les soignants."
    p.level = 1

    # --- SLIDE 5: IMPACT & RESILIENCE ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Stratégie de Résilience"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Anticiper pour protéger :"
    
    p = tf.add_paragraph()
    p.text = "• Alertes précoces conformes aux normes OMS."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Optimisation de la réponse hospitalière lors des pics."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Aide aux politiques de reforestation urbaine."
    p.level = 1

    # --- SLIDE 6: ROADMAP ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Roadmap & Scalabilité"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Expansion du réseau :"
    
    p = tf.add_paragraph()
    p.text = "• Phase 2 : Intégration de capteurs IoT bas coût."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Phase 3 : Notifications SMS pour les zones rurales."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Vision : Extension à l'Afrique Centrale et au Sahel."
    p.level = 1

    # --- SLIDE 7: CONCLUSION ---
    slide = prs.slides.add_slide(slide_layout)
    add_bg(slide, GREEN_ACCENT)
    
    title = slide.shapes.title
    title.text = "SmartAir Cameroon"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Rejoignez la révolution de l'air propre.\nRespirez l'innovation."
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE

    # Save
    output_path = "SmartAir_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Pitch Deck genere avec succes : {output_path}")

if __name__ == "__main__":
    create_pitch_deck()
